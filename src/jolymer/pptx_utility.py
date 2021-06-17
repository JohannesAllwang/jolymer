import pptx
import matplotlib.pyplot as plt
import numpy as np
from io import StringIO
from pptx.util import Inches

class PptxPages:

    def __init__(self, filename, path='../reports/'):
        path = ''
        self.filename = path + filename

    def __enter__(self):
        print(self.filename)
        try:
            # self.prs = pptx.Presentation(self.filename)
            self.prs = pptx.Presentation()
        except:
            print('file not found!')
        #     self.prs = pptx.Presentation()
        return self.prs

    def __exit__(self, exc_class, exc, traceback):
        self.prs.save(self.filename)
        
    def savefig(self, **kwargs):
        pass

class TumTemplate(PptxPages):

    def __init__(self, filename, path='../reports/'):
        PptxPages.__init__(self, filename, path=path)

    def __enter__(self):
        self.prs = pptx.Presentation('jolib/pptx_templates/TUM_Standard.pptx')
        return self.prs

class OpenFile:

    def __init__(self, filename):
        pass
    
def only_figs(filename, plots):
    with PptxPages(filename) as prs:
        for plot in plots:
            x=y=Inches(1)
            width = Inches(5)
            plot()
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            image_stream = StringIO()
            plt.tight_layout()
            plt.savefig('deleteme.jpg')
            pic = slide.shapes.add_picture('deleteme.jpg', x, y, width=width)
            break

def one_par_figs(filename, plot, pars):
    with PptxPages(filename) as prs:
        for par in pars:
            x=y=Inches(1)
            width = Inches(5)
            plot(par)
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            image_stream = StringIO()
            plt.tight_layout()
            plt.savefig('deleteme.jpg')
            pic = slide.shapes.add_picture('deleteme.jpg', x, y, width=width)
            plt.close('all')


def figs_titles_bullets(filename, plots, bpll):
    pass
